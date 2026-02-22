#!/usr/bin/env python3
"""
Convert Office documents to plain text.
Usage: convert-doc.py <input-file> [output-file]

Supports: .docx, .xlsx, .pptx
If no output file is given, prints to stdout.
"""
import sys
import os


def convert_docx(filepath: str) -> str:
    from docx import Document
    doc = Document(filepath)
    lines = []
    for para in doc.paragraphs:
        if para.text.strip():
            lines.append(para.text)

    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append('\t'.join(cells))
        lines.append('')  # blank line between tables

    return '\n'.join(lines)


def convert_xlsx(filepath: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(filepath, read_only=True, data_only=True)
    lines = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f'=== Sheet: {sheet_name} ===')
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else '' for c in row]
            lines.append('\t'.join(cells))
        lines.append('')

    wb.close()
    return '\n'.join(lines)


def convert_pptx(filepath: str) -> str:
    from pptx import Presentation
    prs = Presentation(filepath)
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f'=== Slide {i} ===')
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    lines.append('\t'.join(cells))
        lines.append('')

    return '\n'.join(lines)


def main():
    if len(sys.argv) < 2:
        print('Usage: convert-doc.py <input-file> [output-file]', file=sys.stderr)
        sys.exit(1)

    filepath = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None

    if not os.path.exists(filepath):
        print(f'Error: File not found: {filepath}', file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(filepath)[1].lower()
    converters = {
        '.docx': convert_docx,
        '.xlsx': convert_xlsx,
        '.pptx': convert_pptx,
    }

    converter = converters.get(ext)
    if not converter:
        print(f'Error: Unsupported format: {ext}', file=sys.stderr)
        print(f'Supported: {", ".join(converters.keys())}', file=sys.stderr)
        sys.exit(1)

    try:
        result = converter(filepath)
    except Exception as e:
        print(f'Error converting {filepath}: {e}', file=sys.stderr)
        sys.exit(1)

    if output_file:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(result)
        print(f'Converted to: {output_file}', file=sys.stderr)
    else:
        print(result)


if __name__ == '__main__':
    main()
